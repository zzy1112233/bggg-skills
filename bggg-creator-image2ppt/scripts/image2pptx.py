#!/usr/bin/env python3
"""Build editable PPTX files from a bggg-creator-image2ppt manifest."""

from __future__ import annotations

import argparse
import json
import math
import shutil
import subprocess
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable, Sequence

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Inches, Pt


DEFAULT_SLIDE_WIDTH_IN = 13.333
TRANSPARENT_VALUES = {"none", "transparent", "rgba(0,0,0,0)"}
NAMED_COLORS = {
    "black": "#000000",
    "white": "#ffffff",
    "red": "#ff0000",
    "green": "#008000",
    "blue": "#0000ff",
    "yellow": "#ffff00",
    "gray": "#808080",
    "grey": "#808080",
    "silver": "#c0c0c0",
    "navy": "#000080",
    "orange": "#ffa500",
    "purple": "#800080",
}


@dataclass(frozen=True)
class SlideScale:
    canvas_width: float
    canvas_height: float
    slide_width_in: float
    slide_height_in: float

    @property
    def x_ratio(self) -> float:
        return self.slide_width_in / self.canvas_width

    @property
    def y_ratio(self) -> float:
        return self.slide_height_in / self.canvas_height

    def x(self, value: float | int | None) -> Emu:
        return Inches((float(value or 0)) * self.x_ratio)

    def y(self, value: float | int | None) -> Emu:
        return Inches((float(value or 0)) * self.y_ratio)

    def w(self, value: float | int | None) -> Emu:
        return Inches((float(value or 0)) * self.x_ratio)

    def h(self, value: float | int | None) -> Emu:
        return Inches((float(value or 0)) * self.y_ratio)

    def px_font_to_pt(self, value: float | int) -> float:
        # Use the average axis ratio so text remains visually close when aspect ratio is adjusted.
        return float(value) * ((self.slide_width_in / self.canvas_width) + (self.slide_height_in / self.canvas_height)) * 36


def load_json(path: Path) -> dict[str, Any]:
    with path.open("r", encoding="utf-8") as f:
        return json.load(f)


def write_json(path: Path, payload: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as f:
        json.dump(payload, f, ensure_ascii=False, indent=2)


def resolve_path(path_value: str | None, base_dir: Path) -> Path | None:
    if not path_value:
        return None
    path = Path(path_value).expanduser()
    if not path.is_absolute():
        path = base_dir / path
    return path.resolve()


def parse_color(value: Any, default: str | None = None) -> RGBColor | None:
    if value is None:
        value = default
    if value is None:
        return None
    if isinstance(value, (list, tuple)) and len(value) >= 3:
        return RGBColor(int(value[0]), int(value[1]), int(value[2]))
    text = str(value).strip()
    if not text or text.lower() in TRANSPARENT_VALUES:
        return None
    lowered = text.lower().replace(" ", "")
    if lowered in TRANSPARENT_VALUES:
        return None
    if text.lower() in NAMED_COLORS:
        text = NAMED_COLORS[text.lower()]
    if (text.startswith("rgb(") or text.startswith("rgba(")) and text.endswith(")"):
        body = text[text.find("(") + 1 : -1]
        parts = [part.strip() for part in body.split(",")]
        if len(parts) >= 3:
            return RGBColor(int(float(parts[0])), int(float(parts[1])), int(float(parts[2])))
    if text.startswith("#"):
        text = text[1:]
    if len(text) == 3:
        text = "".join(char * 2 for char in text)
    if len(text) != 6:
        raise ValueError(f"unsupported color value: {value!r}")
    return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))


def shape_type(name: str | None) -> MSO_AUTO_SHAPE_TYPE:
    key = (name or "rect").strip().lower()
    mapping = {
        "rect": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        "rectangle": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        "roundrect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        "roundedrectangle": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        "rounded_rect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        "ellipse": MSO_AUTO_SHAPE_TYPE.OVAL,
        "oval": MSO_AUTO_SHAPE_TYPE.OVAL,
        "circle": MSO_AUTO_SHAPE_TYPE.OVAL,
        "rightarrow": MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        "right_arrow": MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        "notchedrightarrow": MSO_AUTO_SHAPE_TYPE.NOTCHED_RIGHT_ARROW,
        "notched_right_arrow": MSO_AUTO_SHAPE_TYPE.NOTCHED_RIGHT_ARROW,
        "chevron": MSO_AUTO_SHAPE_TYPE.CHEVRON,
        "uparrow": MSO_AUTO_SHAPE_TYPE.UP_ARROW,
        "up_arrow": MSO_AUTO_SHAPE_TYPE.UP_ARROW,
        "heart": MSO_AUTO_SHAPE_TYPE.HEART,
    }
    return mapping.get(key, MSO_AUTO_SHAPE_TYPE.RECTANGLE)


def align_value(value: str | None) -> PP_ALIGN:
    key = (value or "left").strip().lower()
    mapping = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "middle": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    return mapping.get(key, PP_ALIGN.LEFT)


def vertical_anchor(value: str | None) -> MSO_ANCHOR:
    key = (value or "top").strip().lower()
    mapping = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "center": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }
    return mapping.get(key, MSO_ANCHOR.TOP)


def apply_typeface(run: Any, font_name: str) -> None:
    run.font.name = font_name
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = r_pr.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            r_pr.append(node)
        node.set("typeface", font_name)


def image_size(path: Path) -> tuple[int, int]:
    with Image.open(path) as image:
        return image.size


def element_kind(element: dict[str, Any]) -> str:
    return str(element.get("kind") or element.get("type") or "image").strip().lower()


def sorted_elements(elements: Iterable[dict[str, Any]]) -> list[dict[str, Any]]:
    return sorted(
        elements,
        key=lambda item: (
            int(item.get("z", item.get("z_index", item.get("order", 0))) or 0),
            int(item.get("_source_index", 0)),
        ),
    )


def normalize_slides(manifest: dict[str, Any]) -> list[dict[str, Any]]:
    if isinstance(manifest.get("slides"), list):
        slides = manifest["slides"]
    else:
        elements = manifest.get("elements") or manifest.get("layers") or []
        slides = [{"name": manifest.get("name") or "Slide 1", "elements": elements}]
    normalized: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            raise ValueError(f"slide #{slide_index} must be an object")
        elements = slide.get("elements") or slide.get("layers") or []
        for element_index, element in enumerate(elements):
            if isinstance(element, dict):
                element.setdefault("_source_index", element_index)
        normalized.append(
            {
                **slide,
                "name": slide.get("name") or f"Slide {slide_index}",
                "elements": elements,
            }
        )
    return normalized


def scale_for_slide(manifest: dict[str, Any], slide_data: dict[str, Any]) -> SlideScale:
    deck = manifest.get("deck") or manifest.get("canvas") or {}
    canvas_width = float(
        slide_data.get("canvas_width")
        or slide_data.get("width")
        or deck.get("canvas_width")
        or deck.get("width")
        or 1600
    )
    canvas_height = float(
        slide_data.get("canvas_height")
        or slide_data.get("height")
        or deck.get("canvas_height")
        or deck.get("height")
        or 900
    )
    slide_width_in = float(
        slide_data.get("slide_width_in")
        or deck.get("slide_width_in")
        or deck.get("width_in")
        or DEFAULT_SLIDE_WIDTH_IN
    )
    slide_height_in = float(
        slide_data.get("slide_height_in")
        or deck.get("slide_height_in")
        or deck.get("height_in")
        or (slide_width_in * canvas_height / canvas_width)
    )
    if canvas_width <= 0 or canvas_height <= 0:
        raise ValueError("canvas dimensions must be positive")
    return SlideScale(canvas_width, canvas_height, slide_width_in, slide_height_in)


def set_slide_size(presentation: Presentation, scale: SlideScale) -> None:
    presentation.slide_width = Inches(scale.slide_width_in)
    presentation.slide_height = Inches(scale.slide_height_in)


def add_background(slide: Any, element: dict[str, Any], scale: SlideScale, base_dir: Path) -> str:
    file_path = resolve_path(element.get("file") or element.get("path") or element.get("src"), base_dir)
    if file_path:
        if not file_path.exists():
            raise FileNotFoundError(f"background image not found: {file_path}")
        add_image(slide, {**element, "x": 0, "y": 0, "w": scale.canvas_width, "h": scale.canvas_height, "fit": element.get("fit", "cover")}, scale, base_dir)
        return "image-background"

    fill_color = parse_color(element.get("fill") or element.get("color") or element.get("background"), "#ffffff")
    if fill_color:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = fill_color
        return "solid-background"
    return "empty-background"


def add_image(slide: Any, element: dict[str, Any], scale: SlideScale, base_dir: Path) -> None:
    file_path = resolve_path(element.get("file") or element.get("path") or element.get("src"), base_dir)
    if not file_path or not file_path.exists():
        raise FileNotFoundError(f"image not found for element {element.get('name') or element.get('id')}: {file_path}")

    x_px = float(element.get("x", 0))
    y_px = float(element.get("y", 0))
    width_px = element.get("w", element.get("width"))
    height_px = element.get("h", element.get("height"))
    source_width, source_height = image_size(file_path)
    if width_px is None and height_px is None:
        width_px, height_px = source_width, source_height
    elif width_px is None:
        width_px = float(height_px) * source_width / source_height
    elif height_px is None:
        height_px = float(width_px) * source_height / source_width
    width_px = float(width_px)
    height_px = float(height_px)

    fit = str(element.get("fit") or "stretch").lower()
    if fit == "contain":
        image_aspect = source_width / source_height
        box_aspect = width_px / height_px
        if image_aspect > box_aspect:
            final_width = width_px
            final_height = width_px / image_aspect
            final_x = x_px
            final_y = y_px + (height_px - final_height) / 2
        else:
            final_height = height_px
            final_width = height_px * image_aspect
            final_x = x_px + (width_px - final_width) / 2
            final_y = y_px
        slide.shapes.add_picture(
            str(file_path),
            scale.x(final_x),
            scale.y(final_y),
            width=scale.w(final_width),
            height=scale.h(final_height),
        )
        return

    picture = slide.shapes.add_picture(
        str(file_path),
        scale.x(x_px),
        scale.y(y_px),
        width=scale.w(width_px),
        height=scale.h(height_px),
    )
    if fit == "cover":
        image_aspect = source_width / source_height
        box_aspect = width_px / height_px
        if image_aspect > box_aspect:
            crop = max(0.0, min(0.45, 1 - box_aspect / image_aspect)) / 2
            picture.crop_left = crop
            picture.crop_right = crop
        elif image_aspect < box_aspect:
            crop = max(0.0, min(0.45, 1 - image_aspect / box_aspect)) / 2
            picture.crop_top = crop
            picture.crop_bottom = crop


def add_text(slide: Any, element: dict[str, Any], scale: SlideScale) -> None:
    x_px = float(element.get("x", 0))
    y_px = float(element.get("y", 0))
    width_px = float(element.get("w", element.get("width", scale.canvas_width - x_px)))
    height_px = float(element.get("h", element.get("height", 80)))

    shape = slide.shapes.add_textbox(scale.x(x_px), scale.y(y_px), scale.w(width_px), scale.h(height_px))
    shape.name = str(element.get("name") or element.get("id") or "Text")
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = bool(element.get("wrap", True))
    text_frame.margin_left = Emu(0)
    text_frame.margin_right = Emu(0)
    text_frame.margin_top = Emu(0)
    text_frame.margin_bottom = Emu(0)
    text_frame.vertical_anchor = vertical_anchor(element.get("valign") or element.get("vertical_align"))

    raw_text = str(element.get("text", ""))
    lines = raw_text.splitlines() or [""]
    font_size_pt = element.get("font_size_pt")
    if font_size_pt is None:
        font_px = element.get("font_size_px", element.get("font_size", 32))
        font_size_pt = scale.px_font_to_pt(float(font_px))

    color = parse_color(element.get("color"), "#111111")
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.alignment = align_value(element.get("align") or element.get("text_align"))
        paragraph.space_after = Pt(0)
        paragraph.space_before = Pt(0)
        if "line_spacing" in element:
            paragraph.line_spacing = float(element["line_spacing"])
        run = paragraph.add_run()
        run.text = line
        font = run.font
        apply_typeface(run, str(element.get("font_family") or element.get("font") or "Arial"))
        font.size = Pt(float(font_size_pt))
        font.bold = bool(element.get("bold", False))
        font.italic = bool(element.get("italic", False))
        if color:
            font.color.rgb = color
        else:
            font.color.theme_color = MSO_THEME_COLOR.TEXT_1


def add_shape(slide: Any, element: dict[str, Any], scale: SlideScale) -> None:
    shape_name = str(element.get("shape") or element.get("kind") or element.get("type") or "rect")
    if shape_name.lower() == "line":
        x1 = float(element.get("x1", element.get("x", 0)))
        y1 = float(element.get("y1", element.get("y", 0)))
        x2 = float(element.get("x2", x1 + float(element.get("w", element.get("width", 0)))))
        y2 = float(element.get("y2", y1 + float(element.get("h", element.get("height", 0)))))
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, scale.x(x1), scale.y(y1), scale.x(x2), scale.y(y2))
        connector.name = str(element.get("name") or element.get("id") or "Line")
        stroke = parse_color(element.get("stroke") or element.get("color"), "#111111")
        if stroke:
            connector.line.color.rgb = stroke
        connector.line.width = scale.w(float(element.get("stroke_width_px", element.get("stroke_width", 1))))
        return

    x_px = float(element.get("x", 0))
    y_px = float(element.get("y", 0))
    width_px = float(element.get("w", element.get("width", 100)))
    height_px = float(element.get("h", element.get("height", 100)))
    shape = slide.shapes.add_shape(
        shape_type(shape_name),
        scale.x(x_px),
        scale.y(y_px),
        scale.w(width_px),
        scale.h(height_px),
    )
    shape.name = str(element.get("name") or element.get("id") or shape_name)

    fill_color = parse_color(element.get("fill") or element.get("background"))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    stroke_value = element.get("stroke", element.get("border_color"))
    stroke_color = parse_color(stroke_value) if stroke_value is not None else None
    if stroke_color:
        shape.line.color.rgb = stroke_color
        shape.line.width = scale.w(float(element.get("stroke_width_px", element.get("stroke_width", 1))))
    else:
        shape.line.fill.background()


def add_table(slide: Any, element: dict[str, Any], scale: SlideScale) -> None:
    rows = element.get("rows") or element.get("data") or []
    if not rows:
        return
    row_count = len(rows)
    col_count = max(len(row) if isinstance(row, list) else 1 for row in rows)
    x_px = float(element.get("x", 0))
    y_px = float(element.get("y", 0))
    width_px = float(element.get("w", element.get("width", scale.canvas_width - x_px)))
    height_px = float(element.get("h", element.get("height", row_count * 44)))
    frame = slide.shapes.add_table(
        row_count,
        col_count,
        scale.x(x_px),
        scale.y(y_px),
        scale.w(width_px),
        scale.h(height_px),
    )
    frame.name = str(element.get("name") or "Table")
    table = frame.table
    font_size_pt = element.get("font_size_pt")
    if font_size_pt is None:
        font_size_pt = scale.px_font_to_pt(float(element.get("font_size_px", element.get("font_size", 22))))
    text_color = parse_color(element.get("color"), "#111111")
    fill_color = parse_color(element.get("fill"))
    for row_index, row in enumerate(rows):
        values = row if isinstance(row, list) else [row]
        for col_index in range(col_count):
            cell = table.cell(row_index, col_index)
            cell.text = str(values[col_index]) if col_index < len(values) else ""
            if fill_color:
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill_color
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = align_value(element.get("align"))
                for run in paragraph.runs:
                    run.font.size = Pt(float(font_size_pt))
                    apply_typeface(run, str(element.get("font_family") or "Arial"))
                    if text_color:
                        run.font.color.rgb = text_color


def add_element(slide: Any, element: dict[str, Any], scale: SlideScale, base_dir: Path, counters: dict[str, int]) -> None:
    kind = element_kind(element)
    if kind == "background":
        add_background(slide, element, scale, base_dir)
        counters["background"] += 1
    elif kind == "image":
        add_image(slide, element, scale, base_dir)
        counters["image"] += 1
    elif kind == "text":
        add_text(slide, element, scale)
        counters["text"] += 1
    elif kind in {"shape", "rect", "rectangle", "roundrect", "ellipse", "circle", "line"}:
        add_shape(slide, element, scale)
        counters["shape"] += 1
    elif kind == "table":
        add_table(slide, element, scale)
        counters["table"] += 1
    else:
        raise ValueError(f"unsupported element kind: {kind}")


def build_pptx(manifest_path: Path, output_path: Path | None = None, summary_path: Path | None = None) -> dict[str, Any]:
    manifest = load_json(manifest_path)
    base_dir = manifest_path.parent
    slides_data = normalize_slides(manifest)
    if not slides_data:
        raise ValueError("manifest contains no slides")

    first_scale = scale_for_slide(manifest, slides_data[0])
    presentation = Presentation()
    set_slide_size(presentation, first_scale)
    blank_layout = presentation.slide_layouts[6]

    counters = {
        "slides": 0,
        "background": 0,
        "image": 0,
        "text": 0,
        "shape": 0,
        "table": 0,
        "unsupported": 0,
    }
    slide_summaries: list[dict[str, Any]] = []

    for slide_data in slides_data:
        scale = scale_for_slide(manifest, slide_data)
        if not math.isclose(scale.slide_width_in, first_scale.slide_width_in) or not math.isclose(
            scale.slide_height_in, first_scale.slide_height_in
        ):
            raise ValueError("all slides in one PPTX must use the same slide size")
        slide = presentation.slides.add_slide(blank_layout)
        elements = sorted_elements(slide_data.get("elements") or [])
        slide_counters = {"elements": len(elements), "image": 0, "text": 0, "shape": 0, "table": 0, "background": 0}
        for element in elements:
            before = counters.copy()
            add_element(slide, element, scale, base_dir, counters)
            for key in slide_counters:
                if key != "elements":
                    slide_counters[key] += counters[key] - before[key]
        counters["slides"] += 1
        slide_summaries.append({"name": slide_data.get("name"), **slide_counters})

    deck = manifest.get("deck") or {}
    props = presentation.core_properties
    output_stem = output_path.stem if output_path else "image2ppt"
    props.title = str(deck.get("name") or manifest.get("name") or output_stem)
    props.subject = "Generated by bggg-creator-image2ppt"

    output_value = output_path or resolve_path(manifest.get("output") or "output.pptx", base_dir)
    assert output_value is not None
    output_value.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_value)

    summary = {
        "manifest": str(manifest_path.resolve()),
        "output": str(output_value.resolve()),
        "slide_width_in": first_scale.slide_width_in,
        "slide_height_in": first_scale.slide_height_in,
        "canvas_width": first_scale.canvas_width,
        "canvas_height": first_scale.canvas_height,
        "counts": counters,
        "slides": slide_summaries,
    }
    if summary_path:
        write_json(summary_path, summary)
    return summary


def maybe_render_pdf(output_path: Path, out_dir: Path) -> str | None:
    executable = shutil.which("soffice") or shutil.which("libreoffice")
    if not executable:
        return None
    out_dir.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        [executable, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(output_path)],
        check=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
    )
    pdf = out_dir / f"{output_path.stem}.pdf"
    return str(pdf) if pdf.exists() else None


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Build editable PPTX from a bggg-creator-image2ppt manifest.",
        formatter_class=argparse.ArgumentDefaultsHelpFormatter,
    )
    subparsers = parser.add_subparsers(dest="command", required=True)

    build = subparsers.add_parser("build", help="build PPTX from manifest")
    build.add_argument("--manifest", required=True, help="manifest JSON path")
    build.add_argument("--output", help="output PPTX path; defaults to manifest output or output.pptx")
    build.add_argument("--summary", help="summary JSON path")
    build.add_argument("--render-pdf", action="store_true", help="also render a PDF preview with LibreOffice if available")

    return parser


def main(argv: Sequence[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    if args.command == "build":
        manifest_path = Path(args.manifest).expanduser().resolve()
        output_path = Path(args.output).expanduser().resolve() if args.output else None
        summary_path = Path(args.summary).expanduser().resolve() if args.summary else None
        summary = build_pptx(manifest_path, output_path, summary_path)
        if args.render_pdf:
            pdf = maybe_render_pdf(Path(summary["output"]), Path(summary["output"]).parent / "diagnostics")
            summary["pdf_preview"] = pdf
            if summary_path:
                write_json(summary_path, summary)
        print(json.dumps(summary, ensure_ascii=False, indent=2))
        return 0
    raise SystemExit(f"unknown command: {args.command}")


if __name__ == "__main__":
    raise SystemExit(main())
